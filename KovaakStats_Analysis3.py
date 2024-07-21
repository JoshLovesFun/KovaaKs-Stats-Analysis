from os import listdir
import xlwt
from xlwt import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches




# ENTER KOVAAK STATS PATH
path = "C:\Program Files (x86)\Steam\steamapps\common\FPSAimTrainer\FPSAimTrainer\stats"
# THIS FUNCTION CREATES AN ARRAY OF ALL THE FILE NAMES IN THE PATH
files = listdir(path)
files.sort()
# CREATE EXCEL WORKBOOK AND WORKSHEETS
wb = Workbook()
sheet1 = wb.add_sheet('All Stats')
sheet2 = wb.add_sheet('Daily Stats')
sheet3 = wb.add_sheet('Monthly Stats')
# CREATE EXCEL TABLE TITLES
sheet1.write(0, 0, 'Name', xlwt.easyxf('font: bold 1'))
sheet1.write(0, 1, 'Date', xlwt.easyxf('font: bold 1'))
sheet1.write(0, 2, 'Score', xlwt.easyxf('font: bold 1'))
sheet1.write(0, 3, 'Sens', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 0, 'Name', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 1, 'Date', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 2, 'Daily Plays', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 3, 'Ave Score', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 4, 'Max Score', xlwt.easyxf('font: bold 1'))
sheet2.write(0, 5, 'Ave Sens', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 0, 'Name', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 1, 'Date', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 2, 'Monthly Plays', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 3, 'Ave Score', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 4, 'Max Score', xlwt.easyxf('font: bold 1'))
sheet3.write(0, 5, 'Ave Sens', xlwt.easyxf('font: bold 1'))
# ITERATE THROUGH ALL KOVAAK STATS FILES################################################################################
for i in range(0, len(files)):
    # GET TASK NAME FROM FILE NAME
    File_Name = files[i]
    Task_Name = File_Name[0:File_Name.find(" - Challenge - ")]
    # GET TASK DATE FROM FILE NAME
    Date = File_Name[File_Name.find(" - Challenge - ") + 15:File_Name.find(" Stats") - 9]
    # ITERATE THROUGH EVERY LINE OF EACH KOVAAK STATS FILE
    # OPEN STATS FILES
    with open(f"{path}/{files[i]}", newline='\n') as csvfile:
        for ii in csvfile:
            # IF LINE HAS SCORE IN IT
            if "Score" in ii:
                # GET SCORE FROM STATS FILE
                Score = ii[7:].strip()
            if "Horiz Sens" in ii:
                # GET SENS FROM STATS FILE
                Sens = ii[12:].strip()
        # convert valorant to cm/360
        if int(round(float(Sens), 2)) < 1:
            Sens = round(16.33/round(float(Sens), 2), 2)
            Sens = str(Sens)
        # WRITE RESULTS TO TEXT FILE
        sheet1.write(i+1, 0, Task_Name)
        sheet1.write(i+1, 1, Date)
        sheet1.write(i+1, 2, round(float(Score), 2))
        sheet1.write(i+1, 3, round(float(Sens), 2))
# ITERATE THROUGH ALL KOVAAK STATS FILES DAILY##########################################################################
Count = 1
Score_Sum = 0
Sens_Sum = 0
Max_Score = 0
iter = 0
for i in range(0, len(files)):
    # GET TASK NAME FROM FILE NAME
    File_Name = files[i]
    Task_Name = File_Name[0:File_Name.find(" - Challenge - ")]
    # GET TASK DATE FROM FILE NAME
    Date = File_Name[File_Name.find(" - Challenge - ") + 15:File_Name.find(" Stats") - 9]
    # GET DAY
    Day = Date[8:]
    # NEXT STATS
    if i < len(files)-1:
        Future_File_Name = files[i + 1]
        Future_Task_Name = Future_File_Name[0:Future_File_Name.find(" - Challenge - ")]
        Future_Date = Future_File_Name[Future_File_Name.find(" - Challenge - ") + 15:Future_File_Name.find(" Stats") - 9]
        Future_Day = Future_Date[8:]
    # OPEN STATS FILES
    with open(f"{path}/{files[i]}", newline='\n') as csvfile:
        # ITERATE THROUGH EVERY LINE OF EACH KOVAAK STATS FILE
        for ii in csvfile:
            # IF LINE HAS SCORE IN IT
            if "Score" in ii:
                # GET SCORE FROM STATS FILE
                Score = ii[7:].strip()
            if "Horiz Sens" in ii:
                # GET SENS FROM STATS FILE
                Sens = ii[12:].strip()
        # convert valorant to cm/360
        if int(round(float(Sens), 2)) < 1:
            Sens = round(16.33/round(float(Sens), 2), 2)
        # PULL MAX VALUE FROM RANGE
        if int(round(float(Score), 2)) > Max_Score:
            Max_Score = int(round(float(Score), 2))
        # IF FUTURE SCORE IS THE SAME DAY
        if Day == Future_Day and Task_Name == Future_Task_Name and i != len(files)-1:
            Count = Count + 1
            Score_Sum = round(Score_Sum + int(round(float(Score), 2)), 2)
            Sens_Sum = round(Sens_Sum + int(round(float(Sens), 2)), 2)
        # IF FUTURE SCORE IS NOT THE SAME DAY
        else:
            # BE SURE COUNT IS GREATER THEN 1
            if Count > 1:
                Score = round((Score_Sum + int(round(float(Score), 2))) / Count, 2)
                Sens = round((Sens_Sum + int(round(float(Sens), 2))) / Count, 2)
            # WRITE RESULTS TO EXCEL FILE
            sheet2.write(iter + 1, 0, Task_Name)
            sheet2.write(iter + 1, 1, Date)
            sheet2.write(iter + 1, 2, Count)
            sheet2.write(iter + 1, 3, round(float(Score), 2))
            sheet2.write(iter + 1, 4, round(float(Max_Score), 2))
            sheet2.write(iter + 1, 5, round(float(Sens), 2))
            iter = iter + 1
            # RESET VALUES
            Count = 1
            Score_Sum = 0
            Sens_Sum = 0
            Max_Score = 0
# ITERATE THROUGH ALL KOVAAK STATS FILES MONTHLY########################################################################
Count = 1
Score_Sum = 0
Sens_Sum = 0
Max_Score = 0
iter = 0
for i in range(0, len(files)):
    # GET TASK NAME FROM FILE NAME
    File_Name = files[i]
    Task_Name = File_Name[0:File_Name.find(" - Challenge - ")]
    # GET TASK DATE FROM FILE NAME
    Date = File_Name[File_Name.find(" - Challenge - ") + 15:File_Name.find(" Stats") - 12]
    # GET MONTH
    Month = Date[5:7]
    # NEXT STATS
    if i < len(files)-1:
        Future_File_Name = files[i + 1]
        Future_Task_Name = Future_File_Name[0:Future_File_Name.find(" - Challenge - ")]
        Future_Date = Future_File_Name[Future_File_Name.find(" - Challenge - ") + 15:Future_File_Name.find(" Stats") - 12]
        Future_Month = Future_Date[5:7]
    # OPEN STATS FILES
    with open(f"{path}/{files[i]}", newline='\n') as csvfile:
        # ITERATE THROUGH EVERY LINE OF EACH KOVAAK STATS FILE
        for ii in csvfile:
            # IF LINE HAS SCORE IN IT
            if "Score" in ii:
                # GET SCORE FROM STATS FILE
                Score = ii[7:].strip()
            if "Horiz Sens" in ii:
                # GET SENS FROM STATS FILE
                Sens = ii[12:].strip()
        # convert valorant to cm/360
        if int(round(float(Sens), 2)) < 1:
            Sens = round(16.33/round(float(Sens), 2), 2)
        # PULL MAX VALUE FROM RANGE
        if int(round(float(Score), 2)) > Max_Score:
            Max_Score = int(round(float(Score), 2))
        # IF FUTURE SCORE IS THE SAME MONTH
        if Month == Future_Month and Task_Name == Future_Task_Name and i != len(files)-1:
            Count = Count + 1
            Score_Sum = round(Score_Sum + int(round(float(Score), 2)), 2)
            Sens_Sum = round(Sens_Sum + int(round(float(Sens), 2)), 2)
        # IF FUTURE SCORE IS NOT THE SAME MONTH
        else:
            # BE SURE COUNT IS GREATER THEN 1
            if Count > 1:
                Score = round((Score_Sum + int(round(float(Score), 2))) / Count, 2)
                Sens = round((Sens_Sum + int(round(float(Sens), 2))) / Count, 2)
            # WRITE RESULTS TO EXCEL FILE
            sheet3.write(iter + 1, 0, Task_Name)
            sheet3.write(iter + 1, 1, Date)
            sheet3.write(iter + 1, 2, Count)
            sheet3.write(iter + 1, 3, round(float(Score), 2))
            sheet3.write(iter + 1, 4, round(float(Max_Score), 2))
            sheet3.write(iter + 1, 5, round(float(Sens), 2))
            iter = iter + 1
            # RESET VALUES
            Count = 1
            Score_Sum = 0
            Sens_Sum = 0
            Max_Score = 0
# CLOSE EXCEL FILE
wb.save('Kovaak_Stats_Analysis.xls')

# create presentation with 1 slide ------
prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = " Created By python-pptx"
chart_data = ChartData()
chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
chart_data.add_series('Ave Scores',    (24.3, 30.6, 20.2))
chart_data.add_series('Max Scores',    (32.2, 28.4, 34.7))
x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(4)
chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True
chart.series[1].smooth = True

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = " Created By python-pptx"
chart_data = XyChartData()
series_1 = chart_data.add_series('Model 1')
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series('Model 2')
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.series[0].smooth = True
chart.series[1].smooth = True

#https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
prs.save('test1.pptx')
# WRITTEN BY JM contact: jmolvar10@gmail.com
